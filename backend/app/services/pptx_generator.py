"""
PPTX Generator Service
Membuat file presentasi (.pptx) berisi:
  - Background dari assets/background.jpg
  - 4 kartu per slide (2x2 grid)
  - Setiap kartu: nama gudang + tabel stok + 2 foto CCTV
  - Slide dikelompokkan per provinsi
"""

import io
import json
import os
import requests
from collections import OrderedDict
from datetime import date
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from sqlalchemy.orm import Session

from app.models.models import (
    StockCalculation, Photo, Warehouse, WarehousePlant, SAPStock
)

# ── Paths ────────────────────────────────────────────────────────────────────
BASE_DIR = Path(__file__).resolve().parents[2]
ASSETS_DIR = BASE_DIR.parent / "assets"
BACKGROUND_IMG = ASSETS_DIR / "background.jpg"
SLIDES_DIR = BASE_DIR / "storage" / "slides"

# ── Colors ───────────────────────────────────────────────────────────────────
# ── Province geographic order (Barat → Timur, DIY sebelum Jateng) ──────────
PROVINCE_ORDER = [
    "Aceh",
    "Sumatera Utara",
    "Sumatera Barat",
    "Riau",
    "Kepulauan Riau",
    "Jambi",
    "Bengkulu",
    "Sumatera Selatan",
    "Bangka Belitung",
    "Lampung",
    "Banten",
    "DKI Jakarta",
    "Jawa Barat",
    "DI Yogyakarta",
    "Jawa Tengah",
    "Jawa Timur",
    "Bali",
    "Nusa Tenggara Barat",
    "Nusa Tenggara Timur",
    "Kalimantan Barat",
    "Kalimantan Tengah",
    "Kalimantan Selatan",
    "Kalimantan Timur",
    "Kalimantan Utara",
    "Sulawesi Utara",
    "Sulawesi Tengah",
    "Gorontalo",
    "Sulawesi Selatan",
    "Sulawesi Barat",
    "Sulawesi Tenggara",
    "Maluku",
    "Maluku Utara",
    "Papua Barat",
    "Papua",
]

# ── Colors ───────────────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x00, 0x1F, 0x3F)
TURQUOISE = RGBColor(0x40, 0xE0, 0xD0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY = RGBColor(0xD0, 0xD0, 0xD0)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
BLUE_TEXT = RGBColor(0x00, 0x7B, 0xFF)
RED_TEXT = RGBColor(0xCC, 0x00, 0x00)

# ── Layout constants (inches) ───────────────────────────────────────────────
SLIDE_W = 13.333
SLIDE_H = 7.5
COLS, ROWS = 2, 2
CARDS_PER_SLIDE = COLS * ROWS

MARGIN_X = 0.3
MARGIN_TOP = 1.0
MARGIN_BOTTOM = 0.3
GAP_X = 0.3
GAP_Y = 0.25

CARD_W = (SLIDE_W - 2 * MARGIN_X - GAP_X) / COLS
CARD_H = (SLIDE_H - MARGIN_TOP - MARGIN_BOTTOM - GAP_Y) / ROWS

CARD_POSITIONS: List[Tuple[float, float]] = []
for r in range(ROWS):
    for c in range(COLS):
        x = MARGIN_X + c * (CARD_W + GAP_X)
        y = MARGIN_TOP + r * (CARD_H + GAP_Y)
        CARD_POSITIONS.append((x, y))


# ── Helpers ──────────────────────────────────────────────────────────────────

def _fetch_photo_bytes(url: str, timeout: int = 10) -> Optional[bytes]:
    """Download foto dari URL (R2 / mock-storage)."""
    try:
        if url.startswith("/mock-storage/"):
            local = os.path.join(os.getcwd(), url.lstrip("/"))
            if os.path.isfile(local):
                with open(local, "rb") as f:
                    return f.read()
            return None
        resp = requests.get(url, timeout=timeout)
        if resp.status_code == 200:
            return resp.content
    except Exception:
        pass
    return None


def _set_cell_text(cell, text: str, font_size: int = 9, bold: bool = False,
                   color: RGBColor = DARK_TEXT, align: PP_ALIGN = PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _set_cell_fill(cell, color: RGBColor):
    from pptx.oxml.ns import qn
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": f"{color}"})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def _add_rounded_rect(slide, left, top, width, height):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = MID_GRAY
    shape.line.width = Pt(0.75)
    shape.adjustments[0] = 0.04
    return shape


# ── Data fetching ────────────────────────────────────────────────────────────

def _get_stock_for_warehouse(db: Session, tanggal: date, gudang_id: int) -> List[Dict[str, Any]]:
    rows = (
        db.query(StockCalculation)
        .filter(
            StockCalculation.tanggal == tanggal,
            StockCalculation.gudang_id == gudang_id,
        )
        .order_by(StockCalculation.tipe_pupuk)
        .all()
    )
    return [
        {
            "tipe_pupuk": r.tipe_pupuk,
            "stok_fisik": r.stok_fisik,
            "outstanding_so": r.outstanding_so,
            "intransit": r.intransit,
            "stok_admin": r.stok_admin,
        }
        for r in rows
        if any([r.stok_fisik, r.outstanding_so, r.intransit, r.stok_admin])
    ]


def _get_photos_for_warehouse(db: Session, tanggal: date, gudang_id: int) -> List[str]:
    photos = (
        db.query(Photo)
        .filter(Photo.tanggal == tanggal, Photo.gudang_id == gudang_id)
        .order_by(Photo.feed_number)
        .limit(2)
        .all()
    )
    return [p.r2_url for p in photos]


def _get_warehouse_info(db: Session, gudang_id: int) -> Optional[Warehouse]:
    return db.query(Warehouse).filter(Warehouse.id == gudang_id).first()


# ── Grouping by provinsi ────────────────────────────────────────────────────

def _provinsi_sort_key(provinsi: str) -> int:
    """Return sort key berdasarkan urutan geografis (Barat → Timur)."""
    prov_lower = provinsi.strip().lower()
    for i, name in enumerate(PROVINCE_ORDER):
        if name.lower() == prov_lower:
            return i
    return len(PROVINCE_ORDER)


def _group_by_provinsi(db: Session, warehouse_ids: List[int]) -> OrderedDict:
    """Kelompokkan warehouse IDs berdasarkan provinsi, urutkan geografis (Barat → Timur)."""
    temp: Dict[str, List[int]] = {}
    for wid in warehouse_ids:
        w = _get_warehouse_info(db, wid)
        if not w:
            continue
        prov = (w.provinsi or "Lainnya").strip()
        if prov not in temp:
            temp[prov] = []
        temp[prov].append(wid)
    return OrderedDict(
        (prov, temp[prov])
        for prov in sorted(temp.keys(), key=_provinsi_sort_key)
    )


# ── Card rendering ───────────────────────────────────────────────────────────

def _render_card(slide, x: float, y: float, w: float, h: float,
                 warehouse: Warehouse, stocks: List[Dict], photo_urls: List[str]):
    """Render 1 kartu gudang di posisi (x, y) dengan ukuran (w, h)."""

    # ── Card background ──────────────────────────────────────────────────
    _add_rounded_rect(slide, Inches(x), Inches(y), Inches(w), Inches(h))

    pad = 0.12
    inner_x = x + pad
    inner_w = w - 2 * pad
    cursor_y = y + pad

    # ── Title bar ────────────────────────────────────────────────────────
    title_h = 0.30
    title_shape = slide.shapes.add_shape(
        1,
        Inches(inner_x), Inches(cursor_y),
        Inches(inner_w), Inches(title_h),
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()

    # Gudang name + provinsi
    kode = "/".join(sorted([p.kode_plant for p in (warehouse.plants or [])]))
    provinsi = (warehouse.provinsi or "").strip()
    display_name = warehouse.nama_gudang
    if kode:
        display_name += f"  ({kode})"
    if provinsi:
        display_name += f"  —  {provinsi}"

    tf = title_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = display_name
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    tf.margin_left = Inches(0.08)
    tf.margin_top = Inches(0.02)

    cursor_y += title_h + 0.06

    # ── Stock table ──────────────────────────────────────────────────────
    if stocks:
        n_rows = len(stocks) + 1
        n_cols = 5  # Jenis | Fisik | Out.SO | Intransit | Admin
        tbl_w = inner_w
        tbl_h = 0.22 * n_rows

        table_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(inner_x), Inches(cursor_y),
            Inches(tbl_w), Inches(tbl_h),
        )
        table = table_shape.table

        # Column widths
        col_widths = [0.72, 0.55, 0.55, 0.55, 0.63]
        total_ratio = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(inner_w * cw / total_ratio)

        # Header
        headers = ["Jenis", "Stok Fisik", "Out. SO", "Intransit", "Stok Admin"]
        for i, h_text in enumerate(headers):
            cell = table.cell(0, i)
            _set_cell_fill(cell, TURQUOISE)
            _set_cell_text(cell, h_text, font_size=8, bold=True,
                           color=DARK_BLUE,
                           align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

        # Data rows
        for row_idx, stock in enumerate(stocks, start=1):
            values = [
                stock["tipe_pupuk"],
                f'{stock["stok_fisik"]:,.2f}',
                f'{stock["outstanding_so"]:,.2f}',
                f'{stock["intransit"]:,.2f}',
                f'{stock["stok_admin"]:,.2f}',
            ]
            bg = LIGHT_GRAY if row_idx % 2 == 0 else WHITE
            for col_idx, val in enumerate(values):
                cell = table.cell(row_idx, col_idx)
                _set_cell_fill(cell, bg)
                is_admin = col_idx == 4
                if is_admin:
                    admin_val = stock["stok_admin"]
                    color = RED_TEXT if admin_val <= 0 else BLUE_TEXT
                else:
                    color = DARK_TEXT
                _set_cell_text(cell, val, font_size=9, bold=is_admin,
                               color=color,
                               align=PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT)

        cursor_y += tbl_h + 0.08
    else:
        no_data = slide.shapes.add_textbox(
            Inches(inner_x), Inches(cursor_y),
            Inches(inner_w), Inches(0.25),
        )
        tf = no_data.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Belum ada data stok"
        run.font.size = Pt(8)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        run.font.name = "Calibri"
        cursor_y += 0.3

    # ── Photos ───────────────────────────────────────────────────────────
    remaining_h = (y + h - pad) - cursor_y
    photo_w = (inner_w - 0.08) / 2
    photo_h = max(remaining_h, 0.5)

    for p_idx in range(2):
        px = inner_x + p_idx * (photo_w + 0.08)

        if p_idx < len(photo_urls) and photo_urls[p_idx]:
            img_bytes = _fetch_photo_bytes(photo_urls[p_idx])
            if img_bytes:
                try:
                    img_stream = io.BytesIO(img_bytes)
                    slide.shapes.add_picture(
                        img_stream, Inches(px), Inches(cursor_y),
                        Inches(photo_w), Inches(photo_h),
                    )
                    continue
                except Exception:
                    pass

        # Placeholder
        placeholder = slide.shapes.add_shape(
            1, Inches(px), Inches(cursor_y),
            Inches(photo_w), Inches(photo_h),
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
        placeholder.line.color.rgb = MID_GRAY
        placeholder.line.width = Pt(0.5)

        tf = placeholder.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"CCTV {'Depan' if p_idx == 0 else 'Dalam'}\n(Tidak tersedia)"
        run.font.size = Pt(8)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        run.font.name = "Calibri"


# ── Main generator ───────────────────────────────────────────────────────────

def generate_pptx(
    db: Session,
    tanggal: date,
    warehouse_ids: List[int],
) -> Tuple[str, str, int, int]:
    """
    Generate file PPTX dan simpan ke storage/slides/.
    Slide dikelompokkan per provinsi (maks 4 gudang/slide).
    Return: (filename, file_path, gudang_count, slide_count)
    """

    SLIDES_DIR.mkdir(parents=True, exist_ok=True)

    # Group by provinsi, then chunk each group into slides of 4
    prov_groups = _group_by_provinsi(db, warehouse_ids)
    all_chunks: List[List[int]] = []
    for prov, ids in prov_groups.items():
        for i in range(0, len(ids), CARDS_PER_SLIDE):
            all_chunks.append(ids[i:i + CARDS_PER_SLIDE])

    slide_count = len(all_chunks)

    # Create presentation (widescreen 16:9)
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W * 914400))
    prs.slide_height = Emu(int(SLIDE_H * 914400))
    blank_layout = prs.slide_layouts[6]

    # Load background image bytes once
    bg_bytes = None
    if BACKGROUND_IMG.is_file():
        with open(BACKGROUND_IMG, "rb") as f:
            bg_bytes = f.read()

    for chunk in all_chunks:
        slide = prs.slides.add_slide(blank_layout)

        # ── Background ───────────────────────────────────────────────────
        if bg_bytes:
            bg_stream = io.BytesIO(bg_bytes)
            slide.shapes.add_picture(
                bg_stream, Inches(0), Inches(0),
                Inches(SLIDE_W), Inches(SLIDE_H),
            )

        # ── Render cards ─────────────────────────────────────────────────
        for card_idx, gudang_id in enumerate(chunk):
            pos_x, pos_y = CARD_POSITIONS[card_idx]
            warehouse = _get_warehouse_info(db, gudang_id)
            if not warehouse:
                continue

            stocks = _get_stock_for_warehouse(db, tanggal, gudang_id)
            photo_urls = _get_photos_for_warehouse(db, tanggal, gudang_id)

            _render_card(slide, pos_x, pos_y, CARD_W, CARD_H,
                         warehouse, stocks, photo_urls)

    # ── Save ─────────────────────────────────────────────────────────────
    gudang_count = len(warehouse_ids)
    filename = f"Laporan_Stok_{tanggal.isoformat()}_{gudang_count}gudang.pptx"
    file_path = str(SLIDES_DIR / filename)

    prs.save(file_path)

    return filename, file_path, gudang_count, slide_count
